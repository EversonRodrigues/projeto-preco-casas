"""Gera dois .pptx: deck executivo (não-técnico) e deck técnico.

Uso: python "docs/Apresentações/build_decks.py"
Output: docs/Apresentações/{deck_executivo,deck_tecnico}.pptx
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "Apresentações"
METRICS = json.loads((ROOT / "docs" / "metrics.json").read_text(encoding="utf-8"))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x00, 0x7E, 0x8A)
ORANGE = RGBColor(0xE3, 0x7B, 0x00)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_rect(slide, left, top, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def _add_text(
    slide,
    text,
    left,
    top,
    w,
    h,
    *,
    size=18,
    bold=False,
    color=DARK_GRAY,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, items, left, top, w, h, *, size=18, color=DARK_GRAY):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def _title_bar(slide, title, subtitle=None, accent=TEAL):
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
    _add_rect(slide, Inches(0), Inches(0.95), SLIDE_W, Inches(0.08), accent)
    _add_text(
        slide, title, Inches(0.5), Inches(0.18), Inches(12.5), Inches(0.5),
        size=28, bold=True, color=WHITE,
    )
    if subtitle:
        _add_text(
            slide, subtitle, Inches(0.5), Inches(0.6), Inches(12.5), Inches(0.4),
            size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
        )


def _footer(slide, page_num, total, deck_name):
    _add_rect(slide, Inches(0), Inches(7.18), SLIDE_W, Inches(0.32), LIGHT_GRAY)
    _add_text(
        slide, f"Preço de Casas  |  {deck_name}",
        Inches(0.4), Inches(7.22), Inches(8), Inches(0.28),
        size=10, color=DARK_GRAY,
    )
    _add_text(
        slide, f"{page_num} / {total}",
        Inches(11.5), Inches(7.22), Inches(1.4), Inches(0.28),
        size=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT,
    )


def _kpi_card(slide, left, top, w, h, value, label, accent=TEAL):
    _add_rect(slide, left, top, w, h, WHITE)
    _add_rect(slide, left, top, w, Inches(0.08), accent)
    _add_text(
        slide, value, left, top + Inches(0.4), w, Inches(1.2),
        size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, label, left, top + Inches(1.6), w, Inches(0.5),
        size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER,
    )


# ============================================================
# DECK 1 — EXECUTIVO (não-técnico)
# ============================================================

def build_executivo(metrics: dict) -> Presentation:
    prs = _new_prs()
    deck = "Deck Executivo"
    pages = []

    def add(fn):
        pages.append(fn)

    # 1. Capa
    @add
    def slide_capa():
        s = _blank(prs)
        _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
        _add_rect(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.08), ORANGE)
        _add_text(
            s, "Estimador de Preços de Casas",
            Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
            size=44, bold=True, color=WHITE,
        )
        _add_text(
            s, "Como o sistema funciona e que valor entrega",
            Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5),
            size=22, color=RGBColor(0xCC, 0xDD, 0xEE),
        )
        _add_text(
            s, "Apresentação executiva  |  Abril 2026",
            Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
            size=14, color=RGBColor(0xCC, 0xDD, 0xEE),
        )

    # 2. Por que isso importa
    @add
    def slide_problema():
        s = _blank(prs)
        _title_bar(s, "Por que estimar preços de casas?")
        _add_bullets(
            s,
            [
                "Avaliação imobiliária manual depende de corretor — leva dias e tem viés.",
                "Compradores e vendedores precisam de uma referência rápida e isenta.",
                "Bancos e seguradoras precisam de pricing consistente em escala.",
                "Investidores comparam centenas de imóveis — precisam filtrar antes de visitar.",
            ],
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0),
            size=20,
        )
        _add_text(
            s, "Em todos esses casos, um modelo bem treinado é mais barato, mais rápido e mais auditável.",
            Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
            size=16, bold=True, color=ORANGE,
        )

    # 3. O que entregamos
    @add
    def slide_solucao():
        s = _blank(prs)
        _title_bar(s, "O que este projeto entrega")
        items = [
            ("Página de Predição", "Você preenche características da casa e recebe uma estimativa em segundos."),
            ("Página de Insights", "Gráficos mostrando o que mais influencia o preço (qualidade, área, ano de construção, etc.)."),
            ("Página Técnica", "Hiperparâmetros, métricas e versões — para auditoria e checagem técnica."),
            ("API automatizada", "Outros sistemas (CRM, simuladores) podem consultar o modelo sem precisar do app web."),
        ]
        for i, (titulo, desc) in enumerate(items):
            top = Inches(1.5 + i * 1.3)
            _add_rect(s, Inches(0.8), top, Inches(11.7), Inches(1.15), LIGHT_GRAY)
            _add_rect(s, Inches(0.8), top, Inches(0.15), Inches(1.15), TEAL)
            _add_text(s, titulo, Inches(1.1), top + Inches(0.15), Inches(3.5), Inches(0.4), size=18, bold=True, color=NAVY)
            _add_text(s, desc, Inches(1.1), top + Inches(0.55), Inches(11.0), Inches(0.6), size=14, color=DARK_GRAY)

    # 4. Como o usuário interage
    @add
    def slide_como_usar():
        s = _blank(prs)
        _title_bar(s, "Como o usuário interage com o sistema")
        passos = [
            ("1", "Abre a página", "Acessa a URL do app no navegador (Chrome, Edge, qualquer um)."),
            ("2", "Preenche formulário", "Informa o que sabe sobre a casa: qualidade geral, área, ano, vagas, banheiros..."),
            ("3", "Clica em Estimar", "Sistema calcula em menos de 1 segundo."),
            ("4", "Vê o resultado", "Recebe o preço estimado e em qual faixa do mercado de Ames a casa se posiciona."),
        ]
        for i, (num, titulo, desc) in enumerate(passos):
            top = Inches(1.5 + i * 1.3)
            _add_rect(s, Inches(0.8), top, Inches(1.0), Inches(1.0), TEAL)
            _add_text(s, num, Inches(0.8), top + Inches(0.1), Inches(1.0), Inches(0.9), size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _add_text(s, titulo, Inches(2.1), top + Inches(0.1), Inches(10.0), Inches(0.45), size=20, bold=True, color=NAVY)
            _add_text(s, desc, Inches(2.1), top + Inches(0.55), Inches(10.5), Inches(0.5), size=14, color=DARK_GRAY)

    # 5. Resultados — métricas em linguagem de negócio
    @add
    def slide_resultados():
        s = _blank(prs)
        _title_bar(s, "Quão preciso é o sistema?")
        _kpi_card(s, Inches(0.8), Inches(1.6), Inches(3.9), Inches(2.4),
                  f"US$ {metrics['metrics']['mae_usd']:,.0f}".replace(",", "."),
                  "Erro médio em dólares", accent=ORANGE)
        _kpi_card(s, Inches(4.95), Inches(1.6), Inches(3.9), Inches(2.4),
                  f"{metrics['metrics']['r2'] * 100:.0f}%",
                  "Da variação do preço explicada", accent=TEAL)
        _kpi_card(s, Inches(9.1), Inches(1.6), Inches(3.4), Inches(2.4),
                  f"{metrics['metrics']['rmsle']:.3f}",
                  "RMSLE — métrica oficial Kaggle", accent=NAVY)
        _add_text(
            s,
            "Em casas típicas (~US$ 180 mil), o sistema erra em média menos de 10% do valor.",
            Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.5),
            size=18, bold=True, color=NAVY,
        )
        _add_bullets(
            s,
            [
                "Para 88% dos imóveis a estimativa cai dentro de uma faixa razoável para uma decisão informada.",
                "Para os 12% restantes (mansões e casos atípicos) recomenda-se avaliação humana.",
                "A precisão tende a melhorar com mais dados ou features adicionais (localização, fotos).",
            ],
            Inches(0.8), Inches(5.1), Inches(11.7), Inches(2.0),
            size=15,
        )

    # 6. Onde isso é útil
    @add
    def slide_casos_de_uso():
        s = _blank(prs)
        _title_bar(s, "Onde isso pode ser usado")
        items = [
            ("Imobiliária", "Triagem rápida de imóveis recebidos — qual vale a pena visitar."),
            ("Banco / Seguradora", "Pré-avaliação automática para crédito imobiliário e cálculo de prêmios."),
            ("Investidor", "Filtrar centenas de listagens, focar nos imóveis com preço fora da curva."),
            ("Plataforma online", "Mostrar 'preço sugerido' para vendedores que estão criando anúncios."),
            ("Educação / Auditoria", "Material didático sobre como modelos de ML são construídos."),
        ]
        for i, (titulo, desc) in enumerate(items):
            top = Inches(1.5 + i * 1.05)
            _add_rect(s, Inches(0.8), top, Inches(0.15), Inches(0.9), TEAL)
            _add_text(s, titulo, Inches(1.1), top + Inches(0.05), Inches(3.3), Inches(0.5), size=17, bold=True, color=NAVY)
            _add_text(s, desc, Inches(4.5), top + Inches(0.05), Inches(8.3), Inches(0.85), size=14, color=DARK_GRAY)

    # 7. Limitações honestas
    @add
    def slide_limitacoes():
        s = _blank(prs)
        _title_bar(s, "O que o sistema NÃO faz", accent=ORANGE)
        _add_bullets(
            s,
            [
                "Não substitui avaliação humana em casos complexos — é uma referência, não um veredito.",
                "O dataset é da cidade de Ames, Iowa (EUA), 2006-2010. Para outros mercados ou anos, precisa retreinar.",
                "Não inclui o bairro como feature — limitação herdada do pré-processamento original.",
                "Não considera fotos, condições visuais, vista, vizinhança subjetiva.",
                "Não atualiza automaticamente com o mercado — qualidade decai se não retreinar regularmente.",
            ],
            Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0),
            size=18,
        )

    # 8. O que vem na caixa
    @add
    def slide_entregaveis():
        s = _blank(prs)
        _title_bar(s, "O que está pronto para uso")
        items = [
            "Aplicativo web público (link compartilhável)",
            "API para integração com outros sistemas",
            "Modelo treinado, salvo e versionado",
            "Documentação técnica completa (PRD, decisões, plano)",
            "Notebooks com toda análise exploratória",
            "Conjunto de testes automatizados (16 testes)",
            "Repositório no GitHub, com histórico claro de cada milestone",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5), size=20)

    # 9. Próximos passos
    @add
    def slide_proximos():
        s = _blank(prs)
        _title_bar(s, "Próximos passos sugeridos", accent=ORANGE)
        items = [
            ("Curto prazo", "Adicionar bairro como feature — esperamos +5pp em precisão."),
            ("Curto prazo", "Painel de monitoramento: alertar se a qualidade cair."),
            ("Médio prazo", "Coletar dados mais recentes (pós-2010) e adaptar para o mercado brasileiro."),
            ("Médio prazo", "Adicionar fotos da casa via visão computacional."),
            ("Longo prazo", "Treinar modelos por região para capturar dinâmicas locais."),
        ]
        for i, (prazo, desc) in enumerate(items):
            top = Inches(1.5 + i * 1.0)
            cor = ORANGE if "Curto" in prazo else (TEAL if "Médio" in prazo else NAVY)
            _add_rect(s, Inches(0.8), top, Inches(2.2), Inches(0.85), cor)
            _add_text(s, prazo, Inches(0.8), top + Inches(0.15), Inches(2.2), Inches(0.55), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _add_text(s, desc, Inches(3.3), top + Inches(0.15), Inches(9.5), Inches(0.55), size=15, color=DARK_GRAY)

    # 10. Demo
    @add
    def slide_demo():
        s = _blank(prs)
        _title_bar(s, "Como ver o sistema funcionando")
        _add_text(
            s,
            "Demo público (clique para abrir no navegador):",
            Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.5),
            size=18, color=DARK_GRAY,
        )
        _add_text(
            s,
            "https://projeto-preco-casas-rai4onxlzxzwjxvntdxmdv.streamlit.app/",
            Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
            size=18, bold=True, color=TEAL,
        )
        _add_text(
            s,
            "Código-fonte:",
            Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.4),
            size=18, color=DARK_GRAY,
        )
        _add_text(
            s,
            "https://github.com/EversonRodrigues/projeto-preco-casas",
            Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.5),
            size=18, bold=True, color=TEAL,
        )
        _add_rect(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), LIGHT_GRAY)
        _add_text(
            s,
            "Sugestão para a apresentação ao vivo:\n"
            "1) Abrir Predição → preencher 5 campos → mostrar preço.\n"
            "2) Abrir Insights → mostrar quais features mais pesam.\n"
            "3) Mencionar a URL pública para a audiência testar depois.",
            Inches(1.0), Inches(5.1), Inches(11.3), Inches(1.6),
            size=14, color=DARK_GRAY,
        )

    # 11. Q&A
    @add
    def slide_qa():
        s = _blank(prs)
        _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
        _add_text(
            s, "Perguntas?",
            Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0),
            size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        _add_text(
            s, "Para detalhes técnicos, ver Deck Técnico.",
            Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
            size=18, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER,
        )

    total = len(pages)
    for i, fn in enumerate(pages, start=1):
        fn()
    for i, slide in enumerate(prs.slides, start=1):
        if i > 1 and i < total:
            _footer(slide, i, total, deck)
    return prs


# ============================================================
# DECK 2 — TÉCNICO
# ============================================================

def build_tecnico(metrics: dict) -> Presentation:
    prs = _new_prs()
    deck = "Deck Técnico"
    pages = []

    def add(fn):
        pages.append(fn)

    rf_cv = metrics["comparison"]["rf"]["cv_rmsle"]
    xgb_cv = metrics["comparison"]["xgb"]["cv_rmsle"]
    rf_h = metrics["comparison"]["rf"]
    xgb_h = metrics["comparison"]["xgb"]

    # 1. Capa
    @add
    def slide_capa():
        s = _blank(prs)
        _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
        _add_rect(s, Inches(0), Inches(3.0), SLIDE_W, Inches(0.08), TEAL)
        _add_text(s, "Pipeline ML — Preço de Casas", Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0), size=44, bold=True, color=WHITE)
        _add_text(s, "Arquitetura, decisões técnicas e métricas", Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5), size=22, color=RGBColor(0xCC, 0xDD, 0xEE))
        _add_text(s, "Apresentação técnica  |  Abril 2026", Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4), size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    # 2. Sumário
    @add
    def slide_sumario():
        s = _blank(prs)
        _title_bar(s, "Sumário")
        items = [
            "1. Problema, dataset e métrica oficial (RMSLE)",
            "2. Arquitetura do pipeline (src/ → API → Streamlit)",
            "3. Pré-processamento e features (48 numéricas)",
            "4. Estratégia de validação (KFold 5 + holdout 20%)",
            "5. Grid search comparativo: Ridge baseline vs RF vs XGB",
            "6. Hiperparâmetros vencedores e métricas finais",
            "7. Decisões críticas (ADRs em docs/decisions.md)",
            "8. Persistência, API e front-end",
            "9. Cobertura de testes (16 testes pytest + AppTest)",
            "10. Deploy, observabilidade e próximos passos",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=18)

    # 3. Stack
    @add
    def slide_stack():
        s = _blank(prs)
        _title_bar(s, "Stack escolhida e justificativa")
        items = [
            ("Python 3.12", "Escolha pragmática — 3.12 estável, deps ML maduras, compatível com Streamlit Cloud."),
            ("scikit-learn 1.5", "Padrão de fato para grid search + KFold + métricas. Pipeline determinístico."),
            ("XGBoost 2.1", "Estado da arte em tabular regression — ganhou pela menor CV RMSLE."),
            ("FastAPI + Pydantic 2", "Validação de payload nativa; geração automática de OpenAPI."),
            ("Streamlit 1.39", "Front-end com tempo de prototipagem ~10x menor que React+Flask."),
            ("Plotly", "Gráficos interativos no dashboard, sem precisar de servidor de dashboards."),
            ("joblib", "Persistência do modelo — padrão sklearn, suporta NumPy bem."),
            ("pytest", "16 testes cobrindo preprocessing, predict, train e API."),
        ]
        for i, (lib, desc) in enumerate(items):
            top = Inches(1.4 + i * 0.65)
            _add_text(s, lib, Inches(0.8), top, Inches(3.0), Inches(0.5), size=15, bold=True, color=TEAL)
            _add_text(s, desc, Inches(3.9), top, Inches(9.0), Inches(0.5), size=13, color=DARK_GRAY)

    # 4. Arquitetura
    @add
    def slide_arquitetura():
        s = _blank(prs)
        _title_bar(s, "Arquitetura — fluxo de dados")
        diagram = (
            "data/train_3_1.csv  ────►  src.data.load_data\n"
            "                            │\n"
            "                            └──► src.data.split_holdout (20%)\n"
            "                                  │\n"
            "                                  └──► src.preprocessing.prepare_features\n"
            "                                        (48 numéricas, log1p no target)\n"
            "                                        │\n"
            "                              ┌─────────┴─────────┐\n"
            "                              │                   │\n"
            "                       run_grid_search('rf')    run_grid_search('xgb')\n"
            "                       12 combos × 5 folds      24 combos × 5 folds\n"
            "                              │                   │\n"
            "                              └─────────┬─────────┘\n"
            "                                        │\n"
            "                          Vencedor: menor CV RMSLE\n"
            "                                        │\n"
            "                          src.train.save_model\n"
            "                          (best_model.pkl + metadata.json)\n"
            "                                        │\n"
            "                              ┌─────────┴─────────┐\n"
            "                              │                   │\n"
            "                          api/main.py        app/streamlit_app\n"
            "                          POST /predict     com fallback local"
        )
        tb = s.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = diagram
        run.font.size = Pt(13)
        run.font.name = "Consolas"
        run.font.color.rgb = NAVY

    # 5. Dataset
    @add
    def slide_dataset():
        s = _blank(prs)
        _title_bar(s, "Dataset — Ames Housing (Kaggle)")
        items = [
            "Origem: 'House Prices: Advanced Regression Techniques', Kaggle.",
            "1460 imóveis no treino (com SalePrice), 1459 no teste (sem SalePrice).",
            "85 colunas no treino: 50 numéricas + 35 object (categóricas).",
            "Pré-processado pelas Partes 1 e 2 da série didática original.",
            "As 35 colunas object foram descartadas (incluindo Neighborhood).",
            "Sobraram 50 numéricas; tirando Id e SalePrice, ficam 48 features.",
            "Dataset versionado no repo (data/train_3_1.csv, data/test_3_1.csv).",
            "Sem NAs após pré-processamento — confirmado em 01_eda.ipynb.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=16)

    # 6. Pré-processamento
    @add
    def slide_preprocess():
        s = _blank(prs)
        _title_bar(s, "Pré-processamento — src/preprocessing.py")
        items = [
            "select_dtypes(['int64', 'float64']) — apenas colunas numéricas.",
            "Remove Id e SalePrice de X.",
            "y = np.log1p(SalePrice) — alinha com a métrica oficial RMSLE.",
            "FEATURE_NAMES é constante computada no import (single source of truth).",
            "get_feature_defaults() retorna mediana — usado pela UI Streamlit.",
            "Mesma função para treino e inferência — sem skew train/serve.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.0), size=16)
        _add_rect(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.6), LIGHT_GRAY)
        _add_text(
            s,
            "Decisão D-02: object cols descartadas (limitação herdada). "
            "Decisão D-01: log-transform escolhido para alinhar com RMSLE Kaggle. "
            "Decisão D-06: defaults via mediana — apenas para UI, não para treino.",
            Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.3),
            size=13, color=DARK_GRAY,
        )

    # 7. Estratégia de validação
    @add
    def slide_validacao():
        s = _blank(prs)
        _title_bar(s, "Estratégia de validação")
        items = [
            "Holdout 20%: train_test_split(random_state=42) — NUNCA visto pelo grid.",
            "CV interna: KFold(n_splits=5, shuffle=True, random_state=42).",
            "Scoring: neg_root_mean_squared_error no log-target (= RMSLE).",
            "Seleção do modelo final: menor RMSLE no CV (não no holdout — D-07).",
            "Reporte final: métricas do holdout (RMSLE, MAE em USD, R²).",
            "Mesmo split em todos os notebooks (random_state=42) para apples-to-apples.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.0), size=16)
        _add_rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2), LIGHT_GRAY)
        _add_text(
            s,
            "Por que CV para selecionar e holdout para reportar?\n"
            "Usar holdout para selecionar = vazamento. Usar CV para reportar = otimismo.",
            Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0),
            size=13, color=DARK_GRAY,
        )

    # 8. Grid search
    @add
    def slide_grid():
        s = _blank(prs)
        _title_bar(s, "Grid search — espaços de busca")
        # RF
        _add_rect(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(2.6), LIGHT_GRAY)
        _add_text(s, "RandomForest — 12 combos", Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.4), size=18, bold=True, color=NAVY)
        _add_text(
            s,
            "n_estimators: [200, 500]\nmax_depth: [10, 20, None]\nmax_features: ['sqrt', 0.5]\n\n"
            "Total: 2 × 3 × 2 × 5 folds = 60 fits",
            Inches(0.9), Inches(2.0), Inches(5.7), Inches(2.0),
            size=14, color=DARK_GRAY, font="Consolas",
        )
        # XGB
        _add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.6), LIGHT_GRAY)
        _add_text(s, "XGBoost — 24 combos", Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4), size=18, bold=True, color=NAVY)
        _add_text(
            s,
            "learning_rate: [0.05, 0.1]\nmax_depth: [4, 6, 10]\nn_estimators: [300, 800]\ncolsample_bytree: [0.7, 1.0]\n\n"
            "Total: 2 × 3 × 2 × 2 × 5 folds = 120 fits",
            Inches(7.1), Inches(2.0), Inches(5.7), Inches(2.0),
            size=14, color=DARK_GRAY, font="Consolas",
        )
        _add_rect(s, Inches(0.6), Inches(4.4), Inches(12.2), Inches(1.0), TEAL)
        _add_text(
            s, "Original (notebook das Partes 1-3): ~750 fits.  Reduzido aqui: 180 fits — ~4x mais rápido sem sacrificar performance.",
            Inches(0.8), Inches(4.6), Inches(11.8), Inches(0.6),
            size=15, bold=True, color=WHITE,
        )
        _add_text(
            s,
            "n_jobs=-1 no GridSearchCV (paralelismo nas folds × combos). "
            "n_jobs=1 nos estimadores para evitar sobre-paralelismo. "
            "tree_method='hist' no XGB para velocidade em CPU.",
            Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.4),
            size=13, color=DARK_GRAY,
        )

    # 9. Resultados
    @add
    def slide_resultados():
        s = _blank(prs)
        _title_bar(s, "Resultados — comparação no holdout (random_state=42)")
        cols = ["Modelo", "CV RMSLE", "Holdout RMSLE", "MAE (USD)", "R²"]
        rows = [
            ("Ridge (baseline, alpha=1.0)", "—", "ver 02_baseline", "—", "—"),
            ("RandomForest", f"{rf_cv:.4f}", f"{rf_h['rmsle']:.4f}", f"{rf_h['mae_usd']:,.0f}", f"{rf_h['r2']:.4f}"),
            ("XGBoost (vencedor)", f"{xgb_cv:.4f}", f"{xgb_h['rmsle']:.4f}", f"{xgb_h['mae_usd']:,.0f}", f"{xgb_h['r2']:.4f}"),
        ]
        col_widths = [Inches(4.5), Inches(2.0), Inches(2.2), Inches(2.0), Inches(1.5)]
        left = Inches(0.6)
        top = Inches(1.6)
        # header
        x = left
        for i, h in enumerate(cols):
            _add_rect(s, x, top, col_widths[i], Inches(0.5), NAVY)
            _add_text(s, h, x, top + Inches(0.1), col_widths[i], Inches(0.4), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += col_widths[i]
        # rows
        for r_i, row in enumerate(rows):
            y = top + Inches(0.5 + r_i * 0.6)
            x = left
            for c_i, val in enumerate(row):
                bg = WHITE if r_i != 2 else RGBColor(0xFF, 0xF5, 0xE0)
                _add_rect(s, x, y, col_widths[c_i], Inches(0.6), bg)
                bold = (r_i == 2)
                _add_text(s, str(val), x, y + Inches(0.13), col_widths[c_i], Inches(0.45), size=14, color=NAVY, align=PP_ALIGN.CENTER, bold=bold)
                x += col_widths[c_i]

        _add_text(
            s, f"XGBoost venceu por {(rf_cv - xgb_cv) * 100:.1f}pp em CV RMSLE.",
            Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.4),
            size=15, bold=True, color=ORANGE,
        )
        _add_text(
            s,
            f"Critério de aceitação CA-03 do PRD: RMSLE ≤ 0.16 — atendido com folga ({xgb_h['rmsle']:.3f}).",
            Inches(0.8), Inches(4.9), Inches(11.7), Inches(0.5),
            size=14, color=DARK_GRAY,
        )

    # 10. Hiperparâmetros vencedores
    @add
    def slide_hp():
        s = _blank(prs)
        _title_bar(s, f"Hiperparâmetros vencedores — {metrics['winner']}")
        params = metrics["comparison"][metrics["winner"]]["best_params"]
        items = [f"{k}: {v}" for k, v in params.items()]
        _add_bullets(s, items, Inches(0.8), Inches(1.6), Inches(11.7), Inches(3.5), size=20)
        _add_rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.4), LIGHT_GRAY)
        _add_text(
            s,
            "Interpretação:\n"
            "• max_depth=4 + n_estimators=300 = árvores rasas mas muitas — controle de variância.\n"
            "• learning_rate=0.05 + colsample_bytree=0.7 = regularização pelo amostrar features por árvore.",
            Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.2),
            size=13, color=DARK_GRAY,
        )

    # 11. Decisões críticas
    @add
    def slide_decisoes():
        s = _blank(prs)
        _title_bar(s, "Decisões críticas — ADRs em docs/decisions.md")
        decs = [
            ("D-01", "Log-transform no target", "RMSLE oficial Kaggle exige; skew cai de 1.88 para 0.12."),
            ("D-02", "Object cols descartadas", "Limitação herdada do pré-processamento; Neighborhood perdido. Aceito como dívida MVP."),
            ("D-03", "PRD errou em 49 features", "Inspeção mostrou 48; documentado, não corrigido no PRD."),
            ("D-04", "Streamlit chama API + fallback local", "API em dev; src.predict direto em deploy (R-04)."),
            ("D-05", "NaN scrub no metadata.json", "XGBRegressor.get_params['missing']=NaN; convertido a null para JSON válido."),
            ("D-06", "Cache módulo-level _TRAIN_DF", "Side-effect no import documentado; defaults via cópia."),
            ("D-07", "Vencedor por CV, não holdout", "Holdout reservado para reporte; usá-lo para seleção = vazamento."),
        ]
        for i, (id_, titulo, desc) in enumerate(decs):
            top = Inches(1.3 + i * 0.78)
            _add_rect(s, Inches(0.6), top, Inches(0.8), Inches(0.65), TEAL)
            _add_text(s, id_, Inches(0.6), top + Inches(0.13), Inches(0.8), Inches(0.4), size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _add_text(s, titulo, Inches(1.55), top, Inches(4.5), Inches(0.35), size=14, bold=True, color=NAVY)
            _add_text(s, desc, Inches(1.55), top + Inches(0.32), Inches(11.0), Inches(0.4), size=11, color=DARK_GRAY)

    # 12. API
    @add
    def slide_api():
        s = _blank(prs)
        _title_bar(s, "API FastAPI — endpoints")
        rows = [
            ("GET /health", "200 + status, model_loaded boolean (lifespan-cached).", "no auth"),
            ("GET /model-info", "200 + metadata.json completo (lru_cache).", "no auth"),
            ("POST /predict", "HouseFeatures(48 floats) → predicted_price + model_version. 422 se incompleto.", "no auth"),
        ]
        widths = [Inches(3.5), Inches(7.5), Inches(1.7)]
        top = Inches(1.6)
        x = Inches(0.6)
        headers = ["Endpoint", "Comportamento", "Auth"]
        for i, h in enumerate(headers):
            _add_rect(s, x, top, widths[i], Inches(0.5), NAVY)
            _add_text(s, h, x, top + Inches(0.1), widths[i], Inches(0.4), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += widths[i]
        for r_i, row in enumerate(rows):
            y = top + Inches(0.5 + r_i * 0.85)
            x = Inches(0.6)
            for c_i, val in enumerate(row):
                _add_rect(s, x, y, widths[c_i], Inches(0.85), WHITE)
                font = "Consolas" if c_i == 0 else "Calibri"
                _add_text(s, val, x + Inches(0.08), y + Inches(0.18), widths[c_i] - Inches(0.16), Inches(0.65), size=12, color=NAVY if c_i == 0 else DARK_GRAY, font=font)
                x += widths[c_i]
        _add_text(
            s,
            "Latência medida: p50 ~7ms, max ~14ms (TestClient, sem rede). "
            "CORSMiddleware permissivo. Logging estruturado em INFO/WARNING.",
            Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.0),
            size=13, color=DARK_GRAY,
        )

    # 13. Streamlit
    @add
    def slide_streamlit():
        s = _blank(prs)
        _title_bar(s, "Front-end Streamlit — 3 páginas")
        items = [
            ("1_Predicao.py", "Formulário com 14 campos visíveis + 34 defaults (mediana). Botão estima preço e mostra percentil na distribuição do treino."),
            ("2_Insights.py", "6 visualizações Plotly: importância de features, distribuição, preço por OverallQual, scatter holdout, comparação RF vs XGB, correlação top 10."),
            ("3_Tecnico.py", "Hiperparâmetros, métricas, comparação detalhada, versões das libs, links para PRD e código."),
        ]
        for i, (arq, desc) in enumerate(items):
            top = Inches(1.5 + i * 1.4)
            _add_rect(s, Inches(0.8), top, Inches(11.7), Inches(1.2), LIGHT_GRAY)
            _add_rect(s, Inches(0.8), top, Inches(0.15), Inches(1.2), TEAL)
            _add_text(s, arq, Inches(1.1), top + Inches(0.15), Inches(3.5), Inches(0.4), size=15, bold=True, color=NAVY, font="Consolas")
            _add_text(s, desc, Inches(1.1), top + Inches(0.55), Inches(11.0), Inches(0.65), size=13, color=DARK_GRAY)

    # 14. Testes
    @add
    def slide_testes():
        s = _blank(prs)
        _title_bar(s, "Cobertura de testes — 16 testes pytest")
        items = [
            "tests/test_preprocessing.py — 5 testes (shapes, defaults, FEATURE_NAMES sem dup).",
            "tests/test_predict.py — 3 testes (load_model, range, ValueError em chave faltante).",
            "tests/test_train.py — 4 testes (_scrub_nan, evaluate_holdout, save_model + tmp_path).",
            "tests/test_api.py — 4 testes via TestClient (/health, /model-info, /predict 200 e 422).",
            "AppTest — 4 scripts Streamlit validados (0 exceptions; clique em 'Estimar preço' renderiza preço).",
            "Latência /predict medida em test fixture: p50 ~7ms, max ~14ms.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=15)

    # 15. Deploy
    @add
    def slide_deploy():
        s = _blank(prs)
        _title_bar(s, "Deploy")
        items = [
            "GitHub: https://github.com/EversonRodrigues/projeto-preco-casas (público, MIT).",
            "Streamlit Cloud com runtime.txt fixando python-3.12.",
            "URL pública: https://projeto-preco-casas-rai4onxlzxzwjxvntdxmdv.streamlit.app/",
            "11 commits granulares — um por milestone (P-X, M1-M6).",
            "models/best_model.pkl + metadata.json versionados — modelo é parte do produto.",
            "data/*.csv versionados — datasets pré-processados são input canônico.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=15)

    # 16. CA checklist
    @add
    def slide_ca():
        s = _blank(prs)
        _title_bar(s, "Critérios de aceitação — PRD Seção 6")
        cas = [
            ("CA-01", "prepare_features retorna shape (N, 48)", "atendido com 48 (D-03)"),
            ("CA-02", "RF e XGB grid em < 15 min cada", "atendido — poucos minutos cada"),
            ("CA-03", f"RMSLE ≤ 0.16 no holdout 20%", f"atendido — {xgb_h['rmsle']:.3f}"),
            ("CA-04", "best_model.pkl carrega e prediz", "atendido — tests/test_predict.py"),
            ("CA-05", "POST /predict < 500ms", "atendido — p50 ~7ms"),
            ("CA-06", "POST /predict inválido → 422", "atendido — tests/test_api.py"),
            ("CA-07", "Streamlit carrega 3 páginas", "atendido — AppTest 0 exceptions"),
            ("CA-08", "Predicao estima preço em < 2s", "atendido — sub-segundo"),
            ("CA-09", "Insights exibe 6 gráficos", "atendido"),
            ("CA-10", "Deploy público com link no README", "atendido — Streamlit Cloud"),
        ]
        widths = [Inches(1.2), Inches(6.5), Inches(5.0)]
        top = Inches(1.4)
        x = Inches(0.6)
        for i, h in enumerate(["ID", "Critério", "Status"]):
            _add_rect(s, x, top, widths[i], Inches(0.5), NAVY)
            _add_text(s, h, x, top + Inches(0.1), widths[i], Inches(0.4), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += widths[i]
        for r_i, (id_, crit, status) in enumerate(cas):
            y = top + Inches(0.5 + r_i * 0.5)
            x = Inches(0.6)
            for c_i, val in enumerate([id_, crit, status]):
                _add_rect(s, x, y, widths[c_i], Inches(0.5), WHITE if r_i % 2 == 0 else LIGHT_GRAY)
                bold = (c_i == 0)
                _add_text(s, val, x + Inches(0.1), y + Inches(0.13), widths[c_i] - Inches(0.2), Inches(0.4), size=12, color=NAVY if c_i == 0 else DARK_GRAY, bold=bold)
                x += widths[c_i]

    # 17. Próximos passos técnicos
    @add
    def slide_proximos():
        s = _blank(prs)
        _title_bar(s, "Próximos passos técnicos", accent=ORANGE)
        items = [
            "Reincluir Neighborhood via target encoding ou ordinal — esperado +5pp em precisão.",
            "Pinning exato via pip-tools (requirements.lock) para reprodutibilidade bit-a-bit.",
            "GitHub Actions: rodar pytest em PRs + smoke da API + AppTest.",
            "API com auth (header) e rate-limit (slowapi) se for além de demo.",
            "Drift monitoring: log de predições + comparação com distribuição de treino.",
            "SHAP por predição na página Predicao — explica para o usuário final.",
            "Versionamento de modelos via MLflow ou DVC + model registry.",
            "A/B test contra Ridge baseline — confirmar que ganho do XGB é estatisticamente significativo.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=15)

    # 18. Referências
    @add
    def slide_refs():
        s = _blank(prs)
        _title_bar(s, "Referências do repositório")
        items = [
            "docs/PRD.md — design completo (escopo, riscos, CA).",
            "docs/plano_implementacao.md — decomposição em P-1 a M6.5.",
            "docs/decisions.md — 7 ADRs documentando trade-offs.",
            "docs/metrics.json — métricas finais + comparison RF vs XGB.",
            "models/metadata.json — hiperparâmetros, train_date, lib versions.",
            "notebooks/01_eda.ipynb — análise exploratória completa.",
            "notebooks/02_baseline.ipynb — Ridge no mesmo holdout.",
            "notebooks/03_grid_search.ipynb — pipeline refatorado consumindo src/.",
            "notebooks/04_avaliacao_final.ipynb — predições, resíduos, comparação Ridge vs RF vs XGB.",
        ]
        _add_bullets(s, items, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), size=14)

    # 19. Q&A
    @add
    def slide_qa():
        s = _blank(prs)
        _add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
        _add_text(s, "Q & A", Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.0), size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, "github.com/EversonRodrigues/projeto-preco-casas", Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5), size=18, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    total = len(pages)
    for fn in pages:
        fn()
    for i, slide in enumerate(prs.slides, start=1):
        if i > 1 and i < total:
            _footer(slide, i, total, deck)
    return prs


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    exec_path = OUT_DIR / "deck_executivo.pptx"
    tec_path = OUT_DIR / "deck_tecnico.pptx"
    build_executivo(METRICS).save(exec_path)
    build_tecnico(METRICS).save(tec_path)
    print(f"OK -> {exec_path.relative_to(ROOT)}")
    print(f"OK -> {tec_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
